from pptx import Presentation
from copy import deepcopy

# Открытие существующего файла PowerPoint
prs = Presentation('prez.pptx')

# Выбор слайда по его индексу (например, 0 - это первый слайд)
slide = prs.slides[4]
changes = {
    'alltehnicdzkh': 'changes',
    'neisprdzkh': 'changes',
    'percentdzkh': 'changes'
}
# Получение всех текстовых фреймов на слайде
for shape in slide.shapes:
    if shape.has_table:
        table = shape.table
        for row in table.rows:
            for cell in row.cells:
                new_cell = deepcopy(cell)
                # Замена текста в новой ячейке, если он есть в словаре changes
                if str(cell.text).strip() in changes.keys():
                    new_cell.text = changes[str(cell.text).strip()]
                # Замена старой ячейки новой
                cell._element.getparent().replace(cell._element, new_cell._element)

    if hasattr(shape, "text"):
        print(shape.text)  # Вывод текста каждого текстового фрейма

prs.save('modified_presentation.pptx')
